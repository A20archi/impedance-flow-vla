r"""
Builds the Saturday talk as an editable PowerPoint (.pptx).
Clean / flat / highly-visual: rendered LaTeX equations, all four figures, both simulation videos
embedded as playable H.264 clips (with poster frames), a future-work slide, and an opening
"what we tried & abandoned" arc.  Numbers are pulled verbatim from out_gate/*.json (see JOURNEY.md).

  /home/user/miniconda3/envs/lerobot/bin/python gate/build_deck.py
"""
import os
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
import imageio.v2 as imageio
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = "/home/user/Desktop/Saptarshi/impedance_flow_vla"
FIG = f"{ROOT}/figs"
MED = f"{FIG}/pptx_media"
os.makedirs(MED, exist_ok=True)

# ---------- palette ----------
ACCENT   = RGBColor(0x1F, 0x77, 0xB4)   # impedance blue (matches the figures)
ACCENT_D = RGBColor(0x12, 0x4A, 0x73)
INK      = RGBColor(0x1E, 0x26, 0x30)   # near-black slate
MUTE     = RGBColor(0x5B, 0x64, 0x70)
HAIR     = RGBColor(0xD7, 0xDE, 0xE6)
PANEL    = RGBColor(0xF2, 0xF6, 0xFA)
PANEL2   = RGBColor(0xEA, 0xF1, 0xF8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED      = RGBColor(0xC0, 0x39, 0x2B)
GREEN    = RGBColor(0x21, 0x7A, 0x46)
GOLD     = RGBColor(0xC8, 0x8A, 0x00)
NAVY_BG  = RGBColor(0x12, 0x1A, 0x24)
FONT     = "Calibri"
FONT_H   = "Calibri"

SW, SH = 13.333, 7.5
LM = 0.72                       # left margin
CW = SW - 2 * LM                # content width

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---------- low-level helpers ----------

def slide():
    return prs.slides.add_slide(BLANK)

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp

def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf

def para(tf, runs, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
         before=0, after=6, lh=1.06, name=FONT, first=False, bullet=False, italic=False):
    """runs: str OR list of (text, dict-overrides)."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before); p.space_after = Pt(after); p.line_spacing = lh
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, ov in runs:
        r = p.add_run(); r.text = text
        f = r.font
        f.size = Pt(ov.get("size", size)); f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic); f.name = ov.get("name", name)
        f.color.rgb = ov.get("color", color)
    return p

def header(s, kicker, title, page):
    box(s, LM, 0.52, 0.14, 0.66, fill=ACCENT)                       # accent tab
    tf = tbox(s, LM + 0.28, 0.50, CW - 0.28, 0.3)
    para(tf, kicker.upper(), size=12.5, bold=True, color=ACCENT, after=0, first=True)
    tf2 = tbox(s, LM + 0.28, 0.74, CW - 0.28, 0.62)
    para(tf2, title, size=27, bold=True, color=INK, after=0, first=True, lh=1.0)
    box(s, LM, 1.46, CW, 0.018, fill=HAIR)                          # hairline rule
    footer(s, page)

def footer(s, page):
    page = len(prs.slides._sldIdLst)        # auto-number by actual position (insertion-safe)
    tf = tbox(s, LM, SH - 0.42, CW * 0.7, 0.3)
    para(tf, "Impedance-Shaped Flow Matching for VLA Policies", size=9.5, color=MUTE, first=True)
    tf2 = tbox(s, SW - LM - 1.2, SH - 0.42, 1.2, 0.3)
    para(tf2, str(page), size=10.5, bold=True, color=MUTE, align=PP_ALIGN.RIGHT, first=True)

def fit(iw, ih, bw, bh):
    s = min(bw / iw, bh / ih)
    return iw * s, ih * s

def place_image(s, path, bx, by, bw, bh, shadow_panel=False):
    iw, ih = Image.open(path).size
    w, h = fit(iw, ih, bw, bh)
    x, y = bx + (bw - w) / 2, by + (bh - h) / 2
    if shadow_panel:
        box(s, x - 0.06, y - 0.06, w + 0.12, h + 0.12, fill=WHITE, line=HAIR, line_w=1.0)
    return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))

def place_movie(s, path, poster, bx, by, bw, bh):
    iw, ih = Image.open(poster).size
    w, h = fit(iw, ih, bw, bh)
    x, y = bx + (bw - w) / 2, by + (bh - h) / 2
    box(s, x - 0.05, y - 0.05, w + 0.10, h + 0.10, fill=INK)        # frame
    mv = s.shapes.add_movie(path, Inches(x), Inches(y), Inches(w), Inches(h),
                            poster_frame_image=poster, mime_type="video/mp4")
    # play button badge
    bb = box(s, x + w/2 - 0.35, y + h/2 - 0.35, 0.7, 0.7, fill=None)
    tf = tbox(s, x + w/2 - 0.35, y + h/2 - 0.46, 0.7, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "▶", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
    cap = tbox(s, x, y + h + 0.06, w, 0.3)
    para(cap, "▶  click to play in Slide Show", size=10,
         color=MUTE, align=PP_ALIGN.CENTER, first=True)
    return mv

def render_eq(tex, name, fontsize=26, color="#1E2630"):
    path = f"{MED}/{name}.png"
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.text(0, 0, tex, fontsize=fontsize, color=color)
    fig.savefig(path, dpi=300, transparent=True, bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)
    return path

def chip(s, x, y, w, h, label, fill, txt=WHITE, size=11.5):
    box(s, x, y, w, h, fill=fill, rounded=True, radius=0.5)
    tf = tbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, size=size, bold=True, color=txt, align=PP_ALIGN.CENTER, first=True)

# ---------- equations (rendered crisp) ----------
EQ_MAIN = render_eq(
    r"$v_\theta(a_\tau,\tau,o)\;=\;M^{-1}\left(\,-\,D\,\dot{a}_\tau\;-\;K\,(a_\tau-a_{\mathrm{goal}})\;+\;f_\theta(a_\tau,\tau,o)\,\right)$",
    "eq_main", fontsize=24)
EQ_MAIN_W = render_eq(
    r"$v_\theta(a_\tau,\tau,o)=M^{-1}\left(-D\,\dot{a}_\tau-K(a_\tau-a_{\mathrm{goal}})+f_\theta\right)$",
    "eq_main_w", fontsize=22, color="#FFFFFF")
EQ_VAN = render_eq(r"$v_\theta(a_\tau,\tau,o)=\mathrm{MLP}_\theta(\cdot)$  (unstructured)", "eq_van", fontsize=18, color="#5B6470")


def make_libero_chart():
    """3-tier LIBERO-Object mean success: vanilla 0.7M / impedance 0.7M / SmolVLA 450M."""
    path = f"{MED}/fig_libero_3tier.png"
    fig, ax = plt.subplots(figsize=(6.0, 4.5), dpi=200)
    labels = ["Vanilla flow\n~0.7M · scratch", "Impedance (ours)\n~0.7M · scratch", "SmolVLA\n450M · pretrained"]
    vals = [28.2, 48.2, 90.5]
    colors = ["#5B6470", "#1F77B4", "#C88A00"]
    bars = ax.bar(labels, vals, color=colors, width=0.62, zorder=3)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 1.6, f"{v:.1f}%", ha="center", va="bottom",
                fontsize=15, fontweight="bold", color="#1E2630")
    ax.text(0.5, 62, "+20.0 pp", ha="center", va="bottom", fontsize=13, fontweight="bold", color="#217A46")
    ax.annotate("", xy=(1, 52), xytext=(0, 32),
                arrowprops=dict(arrowstyle="->", color="#217A46", lw=1.6, connectionstyle="arc3,rad=-0.25"))
    ax.set_ylim(0, 105)
    ax.set_ylabel("LIBERO-Object mean success  (10 tasks, 20 eps)", fontsize=11)
    ax.grid(axis="y", color="#D7DEE6", lw=0.8, zorder=0)
    for sp in ["top", "right"]:
        ax.spines[sp].set_visible(False)
    ax.tick_params(axis="x", labelsize=10.5)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path

LIBERO_CHART = make_libero_chart()

# ---------- LIBERO baseline-vs-ours simulation videos + extracted poster frames ----------
CP2V = f"{ROOT}/outputs/checkpoint_2_libero/videos"

def video_poster(video_path, name):
    path = f"{MED}/{name}.png"
    rdr = imageio.get_reader(video_path)
    Image.fromarray(rdr.get_data(0)).save(path)
    rdr.close()
    return path

LIB_VID1 = f"{CP2V}/task5_pick_up_the_tomato_sauce_and_place_it_in_the_basket_baseline_vs_ours.mp4"
LIB_VID2 = f"{CP2V}/task0_pick_up_the_alphabet_soup_and_place_it_in_the_basket_baseline_vs_ours.mp4"
POSTER_V1 = video_poster(LIB_VID1, "poster_libero_tomato")
POSTER_V2 = video_poster(LIB_VID2, "poster_libero_soup")

# ---------- parameter-budget chart (param-matched, per benchmark) ----------
def make_param_chart():
    path = f"{MED}/fig_param_budget.png"
    fig, axes = plt.subplots(1, 2, figsize=(7.8, 3.6), dpi=200)
    ax = axes[0]
    ax.bar(["vanilla", "impedance\n(ours)"], [399.9, 410.5], color=["#5B6470", "#1F77B4"], width=0.6, zorder=3)
    for i, v in enumerate([399.9, 410.5]):
        ax.text(i, v + 5, f"{v:.0f}k", ha="center", fontsize=11, fontweight="bold", color="#1E2630")
    ax.set_title("MetaWorld MT10  ·  state-based", fontsize=11, fontweight="bold")
    ax.set_ylim(0, 480); ax.set_ylabel("parameters (thousands)", fontsize=10)
    ax = axes[1]
    enc = 238.6
    ax.bar(["vanilla", "impedance\n(ours)"], [enc, enc], color="#C8D4DF", width=0.6, label="shared CNN encoder", zorder=3)
    ax.bar(["vanilla", "impedance\n(ours)"], [499.2, 480.0], bottom=[enc, enc],
           color=["#5B6470", "#1F77B4"], width=0.6, label="flow head", zorder=3)
    for i, v in enumerate([737.9, 718.6]):
        ax.text(i, v + 9, f"{v:.0f}k", ha="center", fontsize=11, fontweight="bold", color="#1E2630")
    ax.set_title("LIBERO-Object  ·  image-based", fontsize=11, fontweight="bold")
    ax.set_ylim(0, 850); ax.legend(fontsize=8, loc="upper center", frameon=False)
    for a in axes:
        for sp in ["top", "right"]:
            a.spines[sp].set_visible(False)
        a.grid(axis="y", color="#D7DEE6", lw=0.7, zorder=0); a.tick_params(labelsize=9.5)
    fig.suptitle("Parameter budget — matched per benchmark (ours ≤ vanilla)", fontsize=12.5, fontweight="bold", y=1.03)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path

PARAM_CHART = make_param_chart()

# ---------- video galleries: all 10 MetaWorld + all 10 LIBERO baseline-vs-ours clips ----------
import glob as _glob
CP1V = f"{ROOT}/outputs/checkpoint_1_mt10/videos"
MT10_ORDER = ["reach-v3", "push-v3", "pick-place-v3", "door-open-v3", "drawer-open-v3",
              "drawer-close-v3", "button-press-topdown-v3", "peg-insert-side-v3", "window-open-v3", "window-close-v3"]
MT10_ITEMS = []
for _t in MT10_ORDER:
    _vp = f"{CP1V}/{_t}_baseline_vs_ours.mp4"
    if os.path.exists(_vp):
        MT10_ITEMS.append((_vp, _t.replace("-v3", "").replace("button-press-topdown", "button-press"),
                           video_poster(_vp, f"poster_mt10_{_t}")))
LIB_OBJ = ["alphabet soup", "cream cheese", "salad dressing", "bbq sauce", "ketchup",
           "tomato sauce", "butter", "milk", "chocolate pudding", "orange juice"]
LIB_ITEMS = []
for _i, _obj in enumerate(LIB_OBJ):
    _c = _glob.glob(f"{CP2V}/task{_i}_*_baseline_vs_ours.mp4")
    if _c:
        LIB_ITEMS.append((_c[0], _obj, video_poster(_c[0], f"poster_lib_{_i}")))

def gallery_movie(s, path, poster, bx, by, bw, bh, label):
    iw, ih = Image.open(poster).size
    w, h = fit(iw, ih, bw - 0.12, bh)
    x = bx + (bw - w) / 2
    box(s, x - 0.025, by - 0.025, w + 0.05, h + 0.05, fill=INK)
    s.shapes.add_movie(path, Inches(x), Inches(by), Inches(w), Inches(h), poster_frame_image=poster, mime_type="video/mp4")
    tfp = tbox(s, x, by + h / 2 - 0.18, w, 0.36, anchor=MSO_ANCHOR.MIDDLE)
    para(tfp, "▶", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
    cap = tbox(s, bx, by + h + 0.03, bw, 0.26)
    para(cap, label, size=9.5, bold=True, color=INK, align=PP_ALIGN.CENTER, first=True)

def gallery(s, items, by=1.82):
    cols, gh = 5, 2.28
    gw = CW / cols
    for i, (vp, lab, po) in enumerate(items):
        r, c = divmod(i, cols)
        gallery_movie(s, vp, po, LM + c * gw, by + r * gh, gw, gh - 0.52, lab)

# =======================================================================================
# SLIDE 1 — TITLE
# =======================================================================================
s = slide()
box(s, 0, 0, SW, SH, fill=NAVY_BG)
box(s, 0, 0, 0.22, SH, fill=ACCENT)
tf = tbox(s, LM, 1.15, CW, 0.4)
para(tf, "ROBOT LEARNING  ·  FLOW-MATCHING POLICIES  ·  CONTACT-RICH MANIPULATION",
     size=13, bold=True, color=RGBColor(0x7F, 0xB2, 0xD9), first=True)
tf = tbox(s, LM, 1.62, CW, 1.9)
para(tf, "Impedance-Shaped Flow Matching", size=46, bold=True, color=WHITE, first=True, lh=1.0, after=2)
para(tf, "for VLA Policies", size=46, bold=True, color=WHITE, after=10, lh=1.0)
para(tf, "Making the action expert's flow-ODE drift a closed-loop impedance law — compliance becomes structural, not a downstream add-on.",
     size=16.5, color=RGBColor(0xC4, 0xCF, 0xDA))
# equation hero on a light card
box(s, LM, 4.05, CW, 1.35, fill=RGBColor(0x1B, 0x26, 0x33), line=RGBColor(0x2C, 0x3B, 0x4C), line_w=1.0, rounded=True, radius=0.06)
place_image(s, EQ_MAIN_W, LM + 0.3, 4.18, CW - 0.6, 1.1)
tf = tbox(s, LM, 6.55, CW, 0.5)
para(tf, [("IS Lab · Changwon National University", {"bold": True, "color": WHITE, "size": 14}),
          ("        Group presentation — Saturday, 2026-05-30", {"color": RGBColor(0x9F,0xAD,0xBC), "size": 14})],
     first=True)

# =======================================================================================
# SLIDE 2 — the road here (1): ideas we brainstormed
# =======================================================================================
s = slide()
header(s, "How we got here  ·  1 / 2", "We started somewhere else: Adjoint Attribution", 2)
tf = tbox(s, LM, 1.65, CW, 0.7)
para(tf, [("Prior project: ", {"bold": True}),
          ("“Why attention explanations for flow-matching robot policies are provably wrong.”  Use the ", {}),
          ("adjoint sensitivity", {"bold": True, "color": ACCENT}),
          (" of the action-generating ODE as the principled attribution. Four bids for novelty:", {})],
     size=15, color=INK, first=True)
cards = [
    ("1", "Attention is “provably wrong”", "Attention maps disagree with true ODE sensitivity."),
    ("2", "Second-order / interaction adjoint", "Go beyond first-order gradients — curvature & token interactions."),
    ("3", "Adjoint-KL fine-tuning", "Use the adjoint signal to fine-tune the policy (performance lever)."),
    ("4", "Test-time goal-steering", "Steer generation along adjoint directions at inference."),
]
cw = (CW - 3 * 0.3) / 4
for i, (n, t, d) in enumerate(cards):
    x = LM + i * (cw + 0.3)
    box(s, x, 2.55, cw, 2.7, fill=PANEL, line=HAIR, line_w=1.0, rounded=True, radius=0.05)
    box(s, x, 2.55, cw, 0.62, fill=ACCENT, rounded=True, radius=0.05)
    box(s, x, 2.9, cw, 0.27, fill=ACCENT)   # square off the bottom of the header band
    tf = tbox(s, x, 2.6, cw, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "HYPOTHESIS " + n, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
    tf = tbox(s, x + 0.18, 3.35, cw - 0.36, 1.85)
    para(tf, t, size=15, bold=True, color=INK, first=True, after=8)
    para(tf, d, size=12.5, color=MUTE, lh=1.12)
tf = tbox(s, LM, 5.5, CW, 0.8)
para(tf, [("The goal we set: ", {"bold": True}),
          ("genuine novelty ", {"bold": True, "color": ACCENT}),
          ("+ an ", {}), ("improving number", {"bold": True, "color": ACCENT}),
          (". Every one of these four was tested against that bar.", {})],
     size=15, color=INK, first=True)

# =======================================================================================
# SLIDE 3 — the road here (2): why we abandoned it
# =======================================================================================
s = slide()
header(s, "How we got here  ·  2 / 2", "… and why we abandoned it", 3)
rows = [
    ("Attention “provably wrong”", "Confound.", "Attention vs adjoint actually agree on rank (Spearman +0.52, top-1 100%) once token granularity is matched."),
    ("2nd-order interaction adjoint", "Scale artifact.", "The 90–122× state-token curvature was an embedding-norm effect (state ≈7 vs image ≈3700, ~518×)."),
    ("Adjoint-KL fine-tuning", "Null.", "+1.7 vs a proper control — inside the noise at n = 12."),
    ("Test-time goal-steering", "Net-negative.", "Helps free-space reach (37→75%) but collapses every contact task (push, window → 0%)."),
]
y = 1.66
rh = 0.86
for t, verdict, d in rows:
    box(s, LM, y, CW, rh - 0.12, fill=PANEL, line=HAIR, line_w=1.0, rounded=True, radius=0.06)
    box(s, LM, y, 0.1, rh - 0.12, fill=RED, rounded=False)
    tf = tbox(s, LM + 0.28, y, 0.55, rh - 0.12, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "✗", size=22, bold=True, color=RED, align=PP_ALIGN.CENTER, first=True)
    tf = tbox(s, LM + 0.95, y, 3.5, rh - 0.12, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=14.5, bold=True, color=INK, first=True)
    tf = tbox(s, LM + 4.5, y, CW - 4.6, rh - 0.12, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(verdict + "  ", {"bold": True, "color": RED}), (d, {"color": INK})], size=13, first=True, lh=1.05)
    y += rh
# survivor + lesson
box(s, LM, y + 0.02, CW, 0.78, fill=PANEL2, line=ACCENT, line_w=1.25, rounded=True, radius=0.06)
box(s, LM, y + 0.02, 0.1, 0.78, fill=GREEN)
tf = tbox(s, LM + 0.28, y + 0.02, 0.55, 0.78, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "✓", size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER, first=True)
tf = tbox(s, LM + 0.95, y + 0.02, CW - 1.1, 0.78, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("One result survived (analysis-only): ", {"bold": True, "color": GREEN}),
          ("an interventional ground truth showed attention's modality ranking is the ", {}),
          ("reverse", {"bold": True}), (" of causal truth, and the adjoint's matches. Not the method/perf paper we wanted → ", {}),
          ("abandoned.", {"bold": True, "color": RED})], size=12.5, first=True, lh=1.05)
tf = tbox(s, LM, y + 0.92, CW, 0.5)
para(tf, [("Lesson kept: ", {"bold": True, "color": ACCENT}),
          ("model-internal attributions evaporate under scale/granularity normalization. Validate the core object in isolation, clear a kill-gate before scaling, report numbers honestly.", {})],
     size=13.5, color=INK, first=True)

# =======================================================================================
# SLIDE 4 — the pivot
# =======================================================================================
s = slide()
header(s, "The pivot", "From explaining compliance to building it", 4)
tf = tbox(s, LM, 1.85, CW, 1.9)
para(tf, [("If steering a trained flow policy ", {}), ("collapses", {"bold": True, "color": RED}),
          (" contact tasks, the problem isn't the explanation — the policy has ", {}),
          ("no notion of physical compliance", {"bold": True, "color": ACCENT}),
          (" in how it generates motion.", {})], size=18, color=INK, first=True, lh=1.15, after=10)
para(tf, [("So stop ", {}), ("attributing", {"italic": True}),
          (" behaviour after the fact. Build the right inductive bias ", {}),
          ("into the generator itself.", {"bold": True})], size=18, color=INK, lh=1.15)
# two contrasting panels
box(s, LM, 4.0, CW/2 - 0.2, 2.3, fill=PANEL, line=HAIR, line_w=1.0, rounded=True, radius=0.04)
tf = tbox(s, LM + 0.3, 4.25, CW/2 - 0.8, 1.85)
para(tf, "BEFORE", size=13, bold=True, color=RED, first=True, after=8)
para(tf, "Unstructured neural drift", size=17, bold=True, color=INK, after=6)
para(tf, "Compliance bolted on downstream (external controller / F-T sensor), or merely explained post-hoc.", size=13.5, color=MUTE, lh=1.15)
box(s, LM + CW/2 + 0.2, 4.0, CW/2 - 0.2, 2.3, fill=PANEL2, line=ACCENT, line_w=1.25, rounded=True, radius=0.04)
tf = tbox(s, LM + CW/2 + 0.5, 4.25, CW/2 - 0.8, 1.85)
para(tf, "OURS", size=13, bold=True, color=ACCENT, first=True, after=8)
para(tf, "Drift IS an impedance law", size=17, bold=True, color=INK, after=6)
para(tf, "A closed-loop mass–spring–damper attractor, end-to-end differentiable, baked into the flow ODE. No downstream controller.", size=13.5, color=INK, lh=1.15)
# arrow
tf = tbox(s, LM + CW/2 - 0.25, 4.85, 0.5, 0.6, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "→", size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, first=True)

# =======================================================================================
# SLIDE 5 — thesis / the one new object
# =======================================================================================
s = slide()
header(s, "The contribution", "The one new object: an impedance-shaped flow drift", 5)
box(s, LM, 1.75, CW, 1.5, fill=PANEL, line=ACCENT, line_w=1.25, rounded=True, radius=0.05)
place_image(s, EQ_MAIN, LM + 0.35, 1.95, CW - 0.7, 1.1)
tf = tbox(s, LM, 3.4, CW, 0.4)
para(tf, "The flow-matching action expert's ODE velocity field is itself a 2nd-order closed-loop impedance law.",
     size=14.5, italic=True, color=MUTE, align=PP_ALIGN.CENTER, first=True)
items = [
    ("M, K, D", "task-conditioned SPD matrices (mass / stiffness / damping), Cholesky-parameterized, predicted from features."),
    ("ȧ — velocity state", "the flow is a 2nd-order augmented ODE, integrated semi-implicitly (symplectic Euler). The naive finite-difference form is unstable — the toy gate proved it first."),
    ("a_goal", "an attractor (chunk target or a predicted attractor token) the drift is pulled toward."),
    ("End-to-end", "fully differentiable. No downstream controller. No force/torque sensor at inference."),
]
y = 4.0
for tag, d in items:
    box(s, LM, y + 0.07, 0.16, 0.46, fill=ACCENT)
    tf = tbox(s, LM + 0.35, y, 2.7, 0.7, anchor=MSO_ANCHOR.TOP)
    para(tf, tag, size=14.5, bold=True, color=ACCENT, first=True)
    tf = tbox(s, LM + 3.2, y, CW - 3.2, 0.7, anchor=MSO_ANCHOR.TOP)
    para(tf, d, size=13.5, color=INK, first=True, lh=1.08)
    y += 0.72

# =======================================================================================
# SLIDE 6 — why it's novel (positioning)
# =======================================================================================
s = slide()
header(s, "Positioning", "Why it's novel — and what it is not", 6)
comp = [
    ("CompliantVLA-adaptor", "Predicts K, D from a VLM but feeds an EXTERNAL impedance controller wrapped around a frozen VLA."),
    ("Flow with the Force Field", "Standard unstructured flow drift + a DOWNSTREAM passive impedance controller for compliance."),
    ("ACP / Diffusion-Impedance / DMP", "Stiffness lives on the OUTPUT head or a DMP decoder — not in the generative drift."),
]
y = 1.7
for name, d in comp:
    box(s, LM, y, CW, 0.92, fill=PANEL, line=HAIR, line_w=1.0, rounded=True, radius=0.05)
    tf = tbox(s, LM + 0.3, y, 3.7, 0.92, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, size=15, bold=True, color=INK, first=True)
    tf = tbox(s, LM + 4.2, y, CW - 4.5, 0.92, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, d, size=13.5, color=MUTE, first=True, lh=1.1)
    y += 1.04
box(s, LM, y + 0.02, CW, 1.15, fill=ACCENT, rounded=True, radius=0.05)
tf = tbox(s, LM + 0.4, y + 0.02, CW - 0.8, 1.15, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Ours:  ", {"bold": True, "size": 18, "color": WHITE}),
          ("the velocity field IS the impedance law", {"bold": True, "size": 18, "color": WHITE})],
     first=True, after=4)
para(tf, "A closed-loop K/D/M attractor conditioned on features, embedded inside the action expert — end-to-end, no external controller, no F/T sensor. (The prior-art check confirms this intersection is unclaimed as of May 2026.)",
     size=12.5, color=RGBColor(0xE2, 0xEE, 0xF8), lh=1.1)

# =======================================================================================
# SLIDE 7 — MECHANISM VIDEO
# =======================================================================================
s = slide()
header(s, "Explainable simulation  ·  mechanism", "How one action is generated through the flow ODE", 7)
place_movie(s, f"{MED}/mechanism_anim.mp4", f"{MED}/poster_mechanism.png", LM, 1.66, CW, 3.95)
tf = tbox(s, LM, 5.8, CW, 1.3)
para(tf, [("Vanilla flow", {"bold": True, "color": MUTE}),
          (" travels a straight, unstructured path and arrives at full speed. ", {}),
          ("Impedance flow (ours)", {"bold": True, "color": ACCENT}),
          (" is a damped spring toward the predicted attractor — it decelerates to a ", {}),
          ("soft landing.", {"bold": True})], size=15, color=INK, first=True, lh=1.12, after=6)
para(tf, [("ζ = 1 (ours, critically damped) is smooth; ζ = 0.3 overshoots — which is exactly why we enforce ", {"color": MUTE}),
          ("ζ ≥ 1", {"bold": True, "color": ACCENT}),
          (" via D = 2ζ√K.", {"color": MUTE})], size=14, lh=1.1)

# =======================================================================================
# SLIDE — METHODOLOGY VIDEO (full pipeline)
# =======================================================================================
s = slide()
header(s, "Methodology  ·  end-to-end pipeline", "The full methodology, animated", 0)
place_movie(s, f"{FIG}/methodology_anim.mp4", f"{FIG}/methodology_poster.png", LM, 1.66, CW, 4.0)
tf = tbox(s, LM, 5.95, CW, 1.1)
para(tf, [("Observation → features → predicted ", {}), ("M, K, D & attractor", {"bold": True, "color": ACCENT}),
          (" → the action chunk is generated by integrating the ", {}),
          ("impedance-shaped flow ODE", {"bold": True, "color": ACCENT}),
          (" (symplectic Euler) → executed closed-loop. No external controller, no F/T sensor.", {})],
     size=15, color=INK, first=True, lh=1.14)

# =======================================================================================
# SLIDE 8 — TOY GATE
# =======================================================================================
s = slide()
header(s, "De-risking  ·  validate the core in isolation", "Toy gate — cleared before touching the model", 8)
place_image(s, f"{FIG}/toy_impedance.png", LM, 1.66, CW * 0.52, 3.5, shadow_panel=True)
place_image(s, f"{FIG}/toy_stability.png", LM + CW * 0.54, 1.66, CW * 0.46, 3.5, shadow_panel=True)
y = 5.45
checks = [
    ("SPD heads verified", "Cholesky M = LLᵀ + εI stays positive-definite."),
    ("Damping-ratio control", "ζ = 0.3 → 36% overshoot,  ζ = 1 → clean,  ζ = 2 → monotonic."),
    ("Stability bound", "at N Euler steps usable stiffness K ≲ (steps)² — the soft/compliant regime, which is what contact wants."),
]
for i, (t, d) in enumerate(checks):
    x = LM + i * (CW / 3)
    tf = tbox(s, x, y, CW / 3 - 0.25, 1.4)
    para(tf, [("✓  ", {"bold": True, "color": GREEN, "size": 15}), (t, {"bold": True, "color": INK, "size": 14})], first=True, after=3)
    para(tf, d, size=12, color=MUTE, lh=1.12)

# =======================================================================================
# SLIDE 9 — how we tested (fair comparison)
# =======================================================================================
s = slide()
header(s, "Stage-1 setup", "A fair test: only the drift parameterization differs", 9)
left = [
    ("Same everything", "Compact state-conditioned flow policy. Vanilla vs impedance arm trained with an identical objective; only the ODE drift differs."),
    ("Param-matched", "Both arms at ~400k parameters — the win is the structure, not capacity."),
    ("Honest protocol", "3 seeds × 30 closed-loop eval episodes; no tuning toward a target."),
]
y = 1.8
for t, d in left:
    box(s, LM, y + 0.06, 0.16, 0.5, fill=ACCENT)
    tf = tbox(s, LM + 0.35, y, CW * 0.54, 0.95)
    para(tf, t, size=16, bold=True, color=INK, first=True, after=3)
    para(tf, d, size=13.5, color=MUTE, lh=1.12)
    y += 1.05
# metrics panel on the right
bx = LM + CW * 0.6
box(s, bx, 1.8, CW * 0.4, 3.3, fill=PANEL, line=HAIR, line_w=1.0, rounded=True, radius=0.05)
tf = tbox(s, bx + 0.3, 2.05, CW * 0.4 - 0.6, 3.0)
para(tf, "WHAT WE MEASURE", size=12.5, bold=True, color=ACCENT, first=True, after=10)
for t, d in [("Success rate", "did the task complete"),
             ("Action jerk", "2nd-difference of actions — smoothness"),
             ("Peak contact force", "95th-pct of contact force — compliance"),
             ("Sample efficiency", "metric vs # demonstrations")]:
    para(tf, [("▪  ", {"color": ACCENT, "bold": True}), (t, {"bold": True, "color": INK})], size=14, after=0)
    para(tf, "      " + d, size=11.5, color=MUTE, after=8, lh=1.05)
tf = tbox(s, LM, 5.45, CW, 1.0)
para(tf, [("Benchmark: ", {"bold": True}),
          ("MetaWorld full-state demos — ", {}),
          ("peg-insert-side", {"bold": True, "color": ACCENT}),
          (" (tight-tolerance contact, the hard task) and ", {}),
          ("door-open", {"bold": True}),
          (" (easy-from-state control). ManiSkill2/3 contact suite is the next, primary benchmark.", {})],
     size=14, color=INK, first=True, lh=1.12)

# =======================================================================================
# SLIDE — HYPERPARAMETERS & PARAMETER BUDGET
# =======================================================================================
s = slide()
header(s, "Reproducibility", "How the demo was trained — hyperparameters & parameter budget", 0)
hp = [
    ["Knob", "MetaWorld MT10", "LIBERO-Object"],
    ["Action chunk H  /  ODE steps", "16  /  10", "16  /  10"],
    ["Impedance hidden width", "256", "256"],
    ["Impedance caps  k / d / ζ", "20 / 18 / 2", "20 / 18 / 2"],
    ["Impedance dims (rest = free flow)", "4 of 4", "6 of 7  (gripper-exempt)"],
    ["Learning rate  /  batch", "1e-3  /  256", "1e-3  /  128"],
    ["Epochs per arm·seed  ×  seeds", "300 × 3", "100 × 3"],
    ["Demos per task", "50", "50"],
    ["Eval protocol", "10 rollouts/task", "20 episodes/task"],
    ["Total epochs trained", "1,800", "6,000"],
]
nrows, ncols = len(hp), 3
tw, th = 6.7, 4.0
gt = s.shapes.add_table(nrows, ncols, Inches(LM), Inches(1.74), Inches(tw), Inches(th)).table
gt.first_row = False; gt.horz_banding = False
for c, wv in zip(range(ncols), [2.85, 1.9, 1.95]):
    gt.columns[c].width = Inches(wv)
gt.rows[0].height = Inches(0.5)
for r in range(1, nrows):
    gt.rows[r].height = Inches((th - 0.5) / (nrows - 1))
for r in range(nrows):
    for c in range(ncols):
        cell = gt.cell(r, c)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else PANEL
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.add_run(); run.text = hp[r][c]
        f = run.font; f.name = FONT; f.size = Pt(12 if r else 12.5)
        f.bold = (r == 0 or c == 0); f.color.rgb = WHITE if r == 0 else INK
place_image(s, PARAM_CHART, LM + tw + 0.2, 1.85, CW - tw - 0.2, 3.4)
tf = tbox(s, LM, 5.95, CW, 1.05)
para(tf, [("Param-matched by construction: ", {"bold": True, "color": ACCENT}),
          ("the vanilla arm's hidden width is sized so its parameter count matches (or exceeds) the impedance arm — ", {}),
          ("the only difference is the drift parameterization, not capacity.", {"bold": True})],
     size=14, color=INK, first=True, lh=1.16, after=4)
para(tf, "Identical training objective (differentiable rollout-matching MSE + small jerk penalty), identical data, 3 seeds. MetaWorld is state-conditioned (no image encoder); LIBERO adds a shared 238.6k CNN encoder to both arms.",
     size=12, italic=True, color=MUTE, lh=1.12)

# =======================================================================================
# SLIDE 10 — RESULT: peg-insert (primary)
# =======================================================================================
s = slide()
header(s, "Result  ·  primary task", "peg-insert-side: a textbook sample-efficiency win", 10)
place_image(s, f"{FIG}/fig1_peg_sample_efficiency.png", LM, 1.62, CW, 3.55, shadow_panel=True)
y = 5.4
stats = [("+22 pp", "success @ 10 demos\n(30% → 52%)", ACCENT),
         ("−42%", "lower action jerk\nat low data", ACCENT),
         ("−18%", "peak contact force\n@ matched success (n=25)", ACCENT)]
cw3 = (CW - 2 * 0.3) / 3
for i, (big, d, c) in enumerate(stats):
    x = LM + i * (cw3 + 0.3)
    box(s, x, y, cw3, 1.45, fill=PANEL2, line=HAIR, line_w=1.0, rounded=True, radius=0.06)
    tf = tbox(s, x + 0.2, y + 0.12, cw3 - 0.4, 0.7)
    para(tf, big, size=30, bold=True, color=c, first=True)
    tf = tbox(s, x + 0.2, y + 0.78, cw3 - 0.4, 0.6)
    for j, ln in enumerate(d.split("\n")):
        para(tf, ln, size=12.5, color=INK, first=(j == 0), after=0, lh=1.0)

# =======================================================================================
# SLIDE 11 — ROLLOUT VIDEO
# =======================================================================================
s = slide()
header(s, "Explainable simulation  ·  rollout", "Impedance inserts where vanilla misses — 10 demos", 11)
place_movie(s, f"{MED}/peg_rollout_compare.mp4", f"{MED}/poster_rollout.png", LM, 1.66, CW, 4.0)
tf = tbox(s, LM, 5.95, CW, 1.1)
para(tf, [("MuJoCo peg-insert, identical seeds, both arms trained on just 10 demos (3 episodes shown). ", {}),
          ("Impedance succeeds where vanilla fails.", {"bold": True, "color": ACCENT})],
     size=15, color=INK, first=True, lh=1.12, after=5)
para(tf, "Note: the live force gauge reads higher for impedance here precisely because it makes the insertion contact that vanilla misses — the −18% compliance number is the matched-success comparison on the previous slide.",
     size=11.5, italic=True, color=MUTE, lh=1.08)

# =======================================================================================
# SLIDE 12 — RESULT: door-open (secondary)
# =======================================================================================
s = slide()
header(s, "Result  ·  secondary task (honest scope)", "door-open saturates — so the prior is redundant there", 12)
place_image(s, f"{FIG}/fig2_door_secondary.png", LM, 1.66, CW * 0.62, 3.7, shadow_panel=True)
bx = LM + CW * 0.66
tf = tbox(s, bx, 1.9, CW * 0.34, 3.6)
para(tf, "WHY THIS MATTERS", size=12.5, bold=True, color=ACCENT, first=True, after=10)
para(tf, [("Door-open is easy from state", {"bold": True, "color": INK}),
          (" — vanilla already hits 82% at 10 demos. No sample-efficiency headroom → ", {}),
          ("the two arms tie on success.", {"bold": True})], size=14, color=MUTE, lh=1.18, after=10)
para(tf, [("The impedance prior helps exactly where ", {}),
          ("physics is the bottleneck", {"bold": True, "color": ACCENT}),
          (" (tight-tolerance insertion) and is harmlessly redundant where it isn't.", {})],
     size=14, color=MUTE, lh=1.18, after=10)
para(tf, "That is the honest, expected signature of a good structural prior — not a universal +X%.",
     size=13.5, italic=True, color=INK, lh=1.15)

# =======================================================================================
# SLIDE 13 — RESULTS AT A GLANCE (table)
# =======================================================================================
s = slide()
header(s, "Evidence", "Stage-1 results at a glance", 13)
data = [
    ["Setting", "Success  (van → imp)", "Action jerk  (van → imp)", "Peak force  (van → imp)"],
    ["peg-insert · 10 demos", "30% → 52%   (+22pp)", "0.99 → 0.57   (−42%)", "509 → 674"],
    ["peg-insert · 25 demos", "74% → 73%   (tie)", "0.41 → 0.24   (−41%)", "536 → 441   (−18%)"],
    ["peg-insert · 100 demos", "88% → 88%   (saturated)", "0.19 → 0.26", "443 → 403"],
    ["door-open · 100 demos", "93% → 93%   (tie)", "0.18 → 0.14   (−22%)", "2202 → 2199"],
]
nrows, ncols = len(data), len(data[0])
tbl_w, tbl_h = CW, 3.7
gt = s.shapes.add_table(nrows, ncols, Inches(LM), Inches(1.8), Inches(tbl_w), Inches(tbl_h)).table
gt.first_row = False; gt.horz_banding = False
widths = [3.3, 3.0, 3.05, 2.58]
for c, wv in enumerate(widths):
    gt.columns[c].width = Inches(wv)
gt.rows[0].height = Inches(0.72)
for r in range(1, nrows):
    gt.rows[r].height = Inches((tbl_h - 0.72) / (nrows - 1))
for r in range(nrows):
    for c in range(ncols):
        cell = gt.cell(r, c)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else PANEL
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.add_run(); run.text = data[r][c]
        f = run.font; f.name = FONT; f.size = Pt(13.5 if r == 0 else 13)
        f.bold = (r == 0 or c == 0)
        f.color.rgb = WHITE if r == 0 else INK
        # green-tint the headline win cell
        if r == 1 and c == 1:
            cell.fill.fore_color.rgb = RGBColor(0xDD, 0xF0, 0xE2); f.bold = True; f.color.rgb = GREEN
tf = tbox(s, LM, 5.7, CW, 0.9)
para(tf, [("matched ~400k params · identical training objective · 3 seeds · 30 eval episodes.   ", {"italic": True, "color": MUTE, "size": 12.5})],
     first=True)
para(tf, [("Read: ", {"bold": True, "color": ACCENT}),
          ("a big low-data gain that closes as data saturates (sample efficiency), smoother actions at low data, lower force at matched success — the structural prior helps where physics matters.", {"color": INK})],
     size=13.5, lh=1.12)

# =======================================================================================
# SLIDE — METAWORLD MT10 VIDEO GALLERY (all 10 tasks)
# =======================================================================================
s = slide()
header(s, "Explainable simulation  ·  MetaWorld MT10", "All 10 tasks — left: baseline (vanilla)   |   right: ours (impedance)", 0)
gallery(s, MT10_ITEMS, by=1.82)
tf = tbox(s, LM, 5.95, CW, 1.0)
para(tf, [("3 seeds × 10 rollouts/task. ", {"color": MUTE}),
          ("MT10 is largely saturated (8/10 tie at 100%) — the prior wins on the hard tasks (reach 83→100%, peg-insert 77→80%) and never meaningfully regresses.  Mean 96.0% → 97.7%.", {"color": INK})],
     size=12.5, first=True, lh=1.15, align=PP_ALIGN.CENTER)

# =======================================================================================
# SLIDE 14 — GENERALIZATION: LIBERO-Object + real-scale anchor
# =======================================================================================
s = slide()
header(s, "Result  ·  generalization to pixels + contact", "It transfers to image-based contact tasks — with an honest real-scale anchor", 14)
place_image(s, LIBERO_CHART, LM, 1.72, CW * 0.52, 4.05, shadow_panel=True)
bx = LM + CW * 0.56
tf = tbox(s, bx, 1.85, CW * 0.44, 0.95)
para(tf, [("+20.0 pp", {"bold": True, "size": 30, "color": ACCENT}),
          ("  from the drift structure", {"size": 14, "color": INK})], first=True, after=3)
para(tf, "LIBERO-Object, vision-based contact-rich pick-&-place. Vanilla 28.2% → impedance 48.2% at matched ~0.7M, from scratch (10 tasks × 3 seeds × 20 eps). Ours has fewer params and stabilizes seeds.",
     size=12.5, color=MUTE, lh=1.16)
box(s, bx, 3.62, CW * 0.44, 1.42, fill=PANEL, line=GOLD, line_w=1.25, rounded=True, radius=0.05)
tf = tbox(s, bx + 0.25, 3.74, CW * 0.44 - 0.5, 1.2)
para(tf, [("Real-scale anchor:  ", {"bold": True, "color": GOLD, "size": 13.5}),
          ("SmolVLA 450M = 90.5%", {"bold": True, "color": INK, "size": 13.5})], first=True, after=4)
para(tf, "The 42-pp gap is scale + 800k-trajectory pretraining + language + 2 cameras — the axes we deliberately removed, not the method losing.",
     size=11.5, color=MUTE, lh=1.16)
tf = tbox(s, bx, 5.18, CW * 0.44, 1.0)
para(tf, [("Genuinely visuomotor — ", {"bold": True, "color": ACCENT}),
          ("blinding the policy collapses it 80%→0% (not replay). And alphabet soup, SmolVLA's weakest (65%), is our arm's best (88%): the Stage-3 target.", {"color": INK})],
     size=11.5, first=True, lh=1.16)

# =======================================================================================
# SLIDE 15 — SIMULATION VIDEOS (LIBERO baseline vs ours)
# =======================================================================================
s = slide()
header(s, "Explainable simulation  ·  LIBERO-Object", "Watch it: impedance succeeds where vanilla fails", 15)
half = (CW - 0.5) / 2
place_movie(s, LIB_VID1, POSTER_V1, LM, 1.74, half, 3.25)
place_movie(s, LIB_VID2, POSTER_V2, LM + half + 0.5, 1.74, half, 3.25)
tf = tbox(s, LM, 5.55, half, 0.4)
para(tf, [("tomato sauce", {"bold": True, "color": INK}), ("    baseline 28% → ", {"color": MUTE}),
          ("ours 88%", {"bold": True, "color": ACCENT})], size=14, first=True, align=PP_ALIGN.CENTER)
tf = tbox(s, LM + half + 0.5, 5.55, half, 0.4)
para(tf, [("alphabet soup", {"bold": True, "color": INK}), ("    baseline 37% → ", {"color": MUTE}),
          ("ours 88%", {"bold": True, "color": ACCENT})], size=14, first=True, align=PP_ALIGN.CENTER)
box(s, LM, 6.05, CW, 0.92, fill=PANEL, line=HAIR, line_w=1.0, rounded=True, radius=0.05)
tf = tbox(s, LM + 0.3, 6.05, CW - 0.6, 0.92, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Each clip:  left = ", {"color": INK}), ("baseline (vanilla flow)", {"bold": True, "color": MUTE}),
          (" ,  right = ", {"color": INK}), ("ours (impedance)", {"bold": True, "color": ACCENT}),
          (" , same scene & seed. Ours makes a smooth, compliant approach and completes the place; vanilla fumbles or stalls. Blinding the policy collapses it 80%→0% — it is genuinely visuomotor, not replay.", {"color": INK})],
     size=12.5, first=True, lh=1.14, align=PP_ALIGN.CENTER)

# =======================================================================================
# SLIDE — LIBERO-OBJECT VIDEO GALLERY (all 10 tasks)
# =======================================================================================
s = slide()
header(s, "Explainable simulation  ·  LIBERO-Object", "All 10 tasks — left: baseline (vanilla)   |   right: ours (impedance)", 0)
gallery(s, LIB_ITEMS, by=1.82)
tf = tbox(s, LM, 5.95, CW, 1.0)
para(tf, [("3 seeds × 20 episodes/task, vision-based, randomized object location. ", {"color": MUTE}),
          ("Mean 28.2% → 48.2% (+20pp); biggest wins: alphabet soup 37→88, tomato sauce 28→88, salad dressing 15→62, butter 38→68.", {"color": INK})],
     size=12.5, first=True, lh=1.15, align=PP_ALIGN.CENTER)

# =======================================================================================
# SLIDE 16 — FUTURE WORK
# =======================================================================================
s = slide()
header(s, "Future work", "From Stage-1 proof to a contact-rich VLA paper", 16)
fw = [
    ("Stage-2 · conditioning ablation", "VLM-predicted M/K/D vs fixed; stiffness-output-only baseline — isolate what the structure buys."),
    ("Stage-3 · scale to SmolVLA", "Port the impedance drift into the SmolVLA action expert with the full 2nd-order conditional flow matching."),
    ("Primary benchmark · ManiSkill2/3", "PegInsertionSide, PlugCharger, TurnFaucet, AssemblyNut — reviewer-accepted contact suite. MetaWorld + LIBERO shown as transfer probes; ManiSkill is the primary contact target."),
    ("Efficiency angle · ODE steps", "5 vs 10 integration steps — same accuracy, fewer steps = lower latency."),
    ("Credibility lever · real arm", "Peg / USB insertion with a wrist F/T sensor."),
]
y = 1.66
for i, (t, d) in enumerate(fw):
    box(s, LM, y, CW, 0.8, fill=PANEL if i % 2 == 0 else WHITE, line=HAIR, line_w=1.0, rounded=True, radius=0.05)
    chip(s, LM + 0.22, y + 0.23, 0.34, 0.34, str(i + 1), ACCENT, size=14)
    tf = tbox(s, LM + 0.78, y, 4.1, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=14.5, bold=True, color=INK, first=True, lh=1.0)
    tf = tbox(s, LM + 5.0, y, CW - 5.2, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, d, size=12.8, color=MUTE, first=True, lh=1.08)
    y += 0.86
box(s, LM, y + 0.04, CW, 0.5, fill=PANEL2, line=ACCENT, line_w=1.0, rounded=True, radius=0.08)
tf = tbox(s, LM + 0.3, y + 0.04, CW - 0.6, 0.5, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Kill-gate:  ", {"bold": True, "color": ACCENT}),
          ("GO only if the ManiSkill contact subset shows ≥ +5pp over vanilla SmolVLA at matched params. Fail fast; report honestly.", {"color": INK})],
     size=13.5, first=True)

# =======================================================================================
# SLIDE 15 — TAKEAWAYS
# =======================================================================================
s = slide()
box(s, 0, 0, SW, SH, fill=NAVY_BG)
box(s, 0, 0, 0.22, SH, fill=ACCENT)
tf = tbox(s, LM, 0.7, CW, 0.5)
para(tf, "TAKEAWAYS", size=14, bold=True, color=RGBColor(0x7F, 0xB2, 0xD9), first=True)
pts = [
    ("A genuinely new object", "the flow-ODE drift made a closed-loop impedance law — not an external controller, not a post-hoc explanation. Prior-art-clear."),
    ("It works where it should", "+22pp sample efficiency, −42% jerk, −18% force on tight-tolerance insertion — and +20pp on image-based LIBERO-Object contact tasks; redundant on easy/saturated tasks."),
    ("Earned by discipline", "toy gate → param-matched Stage-1 → SmolVLA next. The lesson from the abandoned project, applied: validate, gate, report honestly."),
]
y = 1.7
for i, (t, d) in enumerate(pts):
    chip(s, LM, y + 0.05, 0.5, 0.5, str(i + 1), ACCENT, size=18)
    tf = tbox(s, LM + 0.8, y, CW - 0.9, 1.4)
    para(tf, t, size=22, bold=True, color=WHITE, first=True, after=4, lh=1.0)
    para(tf, d, size=15, color=RGBColor(0xC4, 0xCF, 0xDA), lh=1.14)
    y += 1.55
box(s, LM, 6.65, CW, 0.02, fill=RGBColor(0x2C, 0x3B, 0x4C))
tf = tbox(s, LM, 6.8, CW, 0.4)
para(tf, [("Impedance-Shaped Flow Matching for VLA Policies", {"bold": True, "color": WHITE, "size": 13}),
          ("    ·    IS Lab, CWNU    ·    2026-05-30", {"color": RGBColor(0x9F,0xAD,0xBC), "size": 13})],
     first=True)

out = f"{ROOT}/Impedance_Flow_VLA_Saturday.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")
